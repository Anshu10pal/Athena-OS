import io
import json

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile
from sqlalchemy.orm import Session

from app.agents import prompts
from app.core.llm import chat_json
from app.core.security import get_current_user
from app.db.database import get_db
from app.db.models import User, VaultEntry
from app.services.vector_store import add_memory

router = APIRouter(prefix="/api/presentation", tags=["presentation"])


def _extract_pptx(data: bytes) -> list[str]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        slides.append("\n".join(texts))
    return slides


def _extract_pdf(data: bytes) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return [(page.extract_text() or "") for page in reader.pages]


@router.post("/analyze")
async def analyze(file: UploadFile = File(...), user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    data = await file.read()
    name = (file.filename or "").lower()
    if name.endswith(".pptx"):
        slides = _extract_pptx(data)
    elif name.endswith(".pdf"):
        slides = _extract_pdf(data)
    else:
        raise HTTPException(400, "Upload a .pptx or .pdf file")
    if not any(s.strip() for s in slides):
        raise HTTPException(400, "No readable text found in the deck")

    deck_text = "\n\n".join(f"--- Slide {i + 1} ---\n{s[:1200]}" for i, s in enumerate(slides[:25]))
    result = chat_json(
        [
            {"role": "system", "content": prompts.PRESENTATION},
            {"role": "user", "content": deck_text},
        ],
        fast=False,
    )
    db.add(
        VaultEntry(
            user_id=user.id,
            kind="presentation",
            title=file.filename or "presentation",
            content=json.dumps(result)[:8000],
        )
    )
    db.commit()
    try:
        add_memory(user.id, f"Analyzed presentation '{file.filename}'. Summary: {result.get('executive_summary', '')}", kind="presentation")
    except Exception:
        pass
    return result
